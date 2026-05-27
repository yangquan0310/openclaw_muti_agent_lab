#!/usr/bin/env python3
"""
LayoutImporter — 从现有PPT的slides中提取结构，创建为新的slide_layout。

通过底层XML操作，将源PPT中有代表性的slide结构转换为真正的slide_layout，
添加到目标模板的slide_layouts集合中。

原理：
  1. 从源模板中选择有代表性的slides
  2. 提取slide的XML结构（形状位置、大小、样式）
  3. 基于目标模板的现有slideLayout创建新的layout
  4. 将源slide的非占位符形状转换为占位符，添加到新layout
  5. 更新slideMaster引用和[Content_Types].xml

Usage:
    # 分析源模板，列出有代表性的slides
    python3 LayoutImporter.py --source template_wuhanwenli.pptx --analyze

    # 将指定slides导入为目标模板的layout
    python3 LayoutImporter.py --source template_wuhanwenli.pptx --target template.pptx --slides 0,2,3,12,39

    # 自动导入前N个有代表性的slides
    python3 LayoutImporter.py --source template_wuhanwenli.pptx --target template.pptx --auto --top 5
"""

import argparse
import json
import logging
import shutil
import sys
import tempfile
import zipfile
from copy import deepcopy
from pathlib import Path
from typing import Any, Dict, List, Optional

from lxml import etree

logging.basicConfig(level=logging.INFO, format="%(asctime)s [%(levelname)s] %(message)s")
logger = logging.getLogger(__name__)

# 命名空间
NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"
NS_P = "http://schemas.openxmlformats.org/presentationml/2006/main"
NS_R = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
NS_CT = "http://schemas.openxmlformats.org/package/2006/content-types"
NS_RELS = "http://schemas.openxmlformats.org/package/2006/relationships"


def _ns(tag: str, ns: str = NS_P) -> str:
    return f"{{{ns}}}{tag}"


class LayoutImporter:
    """Layout导入器：从源PPT提取slide结构，创建为新的slide_layout。"""

    def __init__(self, source_path: Path, target_path: Path) -> None:
        self.source_path = source_path
        self.target_path = target_path

    # ========== 分析接口 ==========

    def analyze_source(self) -> List[Dict[str, Any]]:
        """分析源模板，返回有代表性的slides列表。"""
        from pptx import Presentation
        prs = Presentation(str(self.source_path))

        candidates = []
        for i, slide in enumerate(prs.slides):
            non_ph_shapes = [sh for sh in slide.shapes if not sh.is_placeholder]
            shape_types = set()
            texts = []
            for sh in slide.shapes:
                shape_types.add(str(sh.shape_type))
                if sh.has_text_frame and sh.text_frame.text.strip():
                    texts.append(sh.text_frame.text.strip()[:20])

            score = 0
            features = []
            if len(non_ph_shapes) >= 3:
                score += len(non_ph_shapes)
                features.append(f"{len(non_ph_shapes)}个非占位符形状")
            if 'PICTURE (13)' in shape_types:
                score += 3
                features.append("含图片")
            if 'TABLE (19)' in shape_types:
                score += 5
                features.append("含表格")
            if 'GROUP (6)' in shape_types:
                score += 2
                features.append("含组合形状")
            if len([s for s in shape_types if 'AUTO_SHAPE' in s]) > 0:
                score += 1
                features.append("含AutoShape")

            category = self._categorize_slide(texts, shape_types)

            candidates.append({
                "index": i,
                "score": score,
                "category": category,
                "total_shapes": len(slide.shapes),
                "non_placeholder": len(non_ph_shapes),
                "shape_types": list(shape_types),
                "text_preview": texts[:3] if texts else [],
                "features": features,
            })

        candidates.sort(key=lambda x: x["score"], reverse=True)
        return candidates

    def _categorize_slide(self, texts: List[str], shape_types: set) -> str:
        all_text = " ".join(texts).lower()
        if any(t.lower().startswith("part ") for t in texts):
            return "章节分隔"
        if "封面" in all_text or (texts and texts[0].startswith("教育")):
            return "封面"
        if "contents" in all_text or "目录" in all_text:
            return "目录"
        if "思考" in all_text or "习题" in all_text or "讨论" in all_text:
            return "思考题"
        if 'TABLE (19)' in shape_types:
            return "表格页"
        if len(texts) >= 5:
            return "多内容页"
        if len(texts) >= 3:
            return "正文页"
        return "其他"

    # ========== 导入接口 ==========

    def import_slides(self, slide_indices: List[int]) -> Dict[str, Any]:
        """将指定slides导入为目标模板的slide_layout。"""
        backup_path = self.target_path.with_suffix('.pptx.backup')
        shutil.copy2(self.target_path, backup_path)
        logger.info("已备份目标模板: %s", backup_path)

        try:
            result = self._do_import(slide_indices)
            return result
        except Exception as exc:
            shutil.copy2(backup_path, self.target_path)
            logger.error("导入失败，已恢复备份: %s", exc)
            raise

    def _do_import(self, slide_indices: List[int]) -> Dict[str, Any]:
        """执行底层XML导入操作。"""
        # 创建临时目录
        with tempfile.TemporaryDirectory() as tmpdir:
            tmp_path = Path(tmpdir)
            target_extract = tmp_path / "target"
            source_extract = tmp_path / "source"

            # 解压两个模板
            with zipfile.ZipFile(self.target_path, 'r') as z:
                z.extractall(target_extract)
            with zipfile.ZipFile(self.source_path, 'r') as z:
                z.extractall(source_extract)

            # 获取当前layout数量
            master_xml_path = target_extract / "ppt" / "slideMasters" / "slideMaster1.xml"
            master_root = etree.parse(str(master_xml_path)).getroot()
            sldLayoutIdLst = master_root.find(_ns("sldLayoutIdLst"))
            original_layout_count = len(sldLayoutIdLst)

            imported = []
            next_layout_num = self._get_next_layout_number(target_extract)

            for src_idx in slide_indices:
                layout_info = self._import_single_slide(
                    source_extract, target_extract, src_idx, next_layout_num
                )
                if layout_info:
                    imported.append(layout_info)
                    next_layout_num += 1

            # 重新打包
            self._repack_pptx(target_extract, self.target_path)

            return {
                "success": True,
                "original_layout_count": original_layout_count,
                "new_layout_count": original_layout_count + len(imported),
                "imported": imported,
            }

    def _get_next_layout_number(self, extract_dir: Path) -> int:
        """获取下一个可用的slideLayout编号。"""
        layout_dir = extract_dir / "ppt" / "slideLayouts"
        existing = []
        for f in layout_dir.glob("slideLayout*.xml"):
            if f.stem.startswith("slideLayout"):
                try:
                    num = int(f.stem.replace("slideLayout", ""))
                    existing.append(num)
                except ValueError:
                    pass
        return max(existing) + 1 if existing else 1

    def _import_single_slide(
        self,
        source_extract: Path,
        target_extract: Path,
        src_slide_idx: int,
        layout_num: int,
    ) -> Optional[Dict[str, Any]]:
        """将单个slide导入为新的slideLayout。"""
        # 读取源slide的XML
        slide_xml_path = source_extract / "ppt" / "slides" / f"slide{src_slide_idx + 1}.xml"
        if not slide_xml_path.exists():
            logger.warning("Slide %d XML不存在", src_slide_idx)
            return None

        slide_root = etree.parse(str(slide_xml_path)).getroot()
        slide_spTree = slide_root.find(_ns("cSld")).find(_ns("spTree"))

        # 读取源slide的关系文件，提取图片映射
        slide_rels_path = source_extract / "ppt" / "slides" / "_rels" / f"slide{src_slide_idx + 1}.xml.rels"
        image_map = {}
        if slide_rels_path.exists():
            rels_root = etree.parse(str(slide_rels_path)).getroot()
            for rel in rels_root:
                if "image" in rel.get("Type", ""):
                    image_map[rel.get("Id")] = rel.get("Target")

        # 基于目标模板的空白layout创建新layout
        target_layout_dir = target_extract / "ppt" / "slideLayouts"
        base_layout_path = target_layout_dir / "slideLayout6.xml"  # 空白layout
        if not base_layout_path.exists():
            base_layout_path = target_layout_dir / "slideLayout1.xml"

        # 复制base layout
        new_layout_path = target_layout_dir / f"slideLayout{layout_num}.xml"
        shutil.copy2(base_layout_path, new_layout_path)

        # 修改新layout的XML
        layout_root = etree.parse(str(new_layout_path)).getroot()
        layout_spTree = layout_root.find(_ns("cSld")).find(_ns("spTree"))

        # 清空现有形状（保留占位符结构）
        self._clean_layout_shapes(layout_spTree)

        # 从源slide复制非占位符形状到新layout
        copied_count = 0
        shape_idx = 0
        for child in list(slide_spTree):
            tag = child.tag.split("}")[-1] if "}" in child.tag else child.tag

            # 跳过占位符和组容器属性
            if tag in ("nvGrpSpPr", "grpSpPr"):
                continue

            # 检查是否是占位符
            ph = child.find(f".//{{{NS_P}}}ph")
            if ph is not None:
                continue  # 跳过占位符

            # 深拷贝形状
            new_shape = deepcopy(child)

            # 递归处理形状（处理组合形状、转换文本框为占位符、清空文本）
            shape_idx = self._process_shape(new_shape, shape_idx)

            # 添加到layout的spTree
            layout_spTree.append(new_shape)
            copied_count += 1

        # 获取分类
        from pptx import Presentation
        src_prs = Presentation(str(self.source_path))
        texts = []
        if src_slide_idx < len(src_prs.slides):
            for sh in src_prs.slides[src_slide_idx].shapes:
                if sh.has_text_frame and sh.text_frame.text.strip():
                    texts.append(sh.text_frame.text.strip()[:20])
        category = self._categorize_slide(texts, set())

        # 设置layout名称
        cSld = layout_root.find(_ns("cSld"))
        if cSld is not None:
            cSld.set("name", f"武汉文理_{category}_{src_slide_idx}")

        # 保存修改后的layout
        etree.ElementTree(layout_root).write(
            str(new_layout_path), xml_declaration=True, encoding="UTF-8", standalone=True
        )

        # 复制layout的关系文件
        base_rels_path = target_extract / "ppt" / "slideLayouts" / "_rels" / "slideLayout6.xml.rels"
        if not base_rels_path.exists():
            base_rels_path = target_extract / "ppt" / "slideLayouts" / "_rels" / "slideLayout1.xml.rels"
        new_rels_path = target_extract / "ppt" / "slideLayouts" / "_rels" / f"slideLayout{layout_num}.xml.rels"
        shutil.copy2(base_rels_path, new_rels_path)

        # 更新[Content_Types].xml
        self._update_content_types(target_extract, layout_num)

        # 更新slideMaster的sldLayoutIdLst
        layout_index = self._update_slide_master(target_extract, layout_num)

        return {
            "source_slide": src_slide_idx,
            "layout_file": f"slideLayout{layout_num}.xml",
            "layout_index": layout_index,
            "name": f"武汉文理_{category}_{src_slide_idx}",
            "category": category,
            "copied_shapes": copied_count,
        }

    def _clean_layout_shapes(self, spTree: etree.Element) -> None:
        """清空layout中的占位符文本（保留结构）。"""
        for shape in list(spTree):
            tag = shape.tag.split("}")[-1] if "}" in shape.tag else shape.tag
            if tag not in ("nvGrpSpPr", "grpSpPr"):
                self._clear_shape_text(shape)

    def _process_shape(self, shape_el: etree.Element, shape_idx: int) -> int:
        """递归处理形状：转换文本框为占位符、清空文本、递归处理组合形状。

        返回下一个可用的占位符索引。
        """
        tag = shape_el.tag.split("}")[-1] if "}" in shape_el.tag else shape_el.tag

        if tag == "grpSp":
            # 递归处理组合形状内部的子形状
            for sub in list(shape_el):
                sub_tag = sub.tag.split("}")[-1] if "}" in sub.tag else sub.tag
                if sub_tag in ("nvGrpSpPr", "grpSpPr"):
                    continue
                # 检查子形状是否已经是占位符
                sub_ph = sub.find(f".//{{{NS_P}}}ph")
                if sub_ph is not None:
                    continue
                shape_idx = self._process_shape(sub, shape_idx)
            return shape_idx

        if tag in ("sp", "cxnSp"):
            # 只有包含实际文本的形状才转为占位符
            if self._has_text_content(shape_el):
                self._convert_to_placeholder(shape_el, shape_idx)
                shape_idx += 1
            # 清空文本内容
            self._clear_shape_text(shape_el)
            return shape_idx

        # 其他类型（pic等）不做处理
        return shape_idx

    def _has_text_content(self, shape_el: etree.Element) -> bool:
        """检查形状是否包含实际文本内容。"""
        for txBody in shape_el.findall(f".//{{{NS_P}}}txBody"):
            for p in txBody.findall(f"{{{NS_A}}}p"):
                # 检查 run 中的文本
                for r in p.findall(f"{{{NS_A}}}r"):
                    t = r.find(f"{{{NS_A}}}t")
                    if t is not None and t.text and t.text.strip():
                        return True
                # 检查段落级别的直接文本
                for t in p.findall(f"{{{NS_A}}}t"):
                    if t.text and t.text.strip():
                        return True
        return False

    def _clear_shape_text(self, shape_el: etree.Element) -> None:
        """清空形状中的所有文本。"""
        for txBody in shape_el.findall(f".//{{{NS_P}}}txBody"):
            for p in txBody.findall(f"{{{NS_A}}}p"):
                for r in p.findall(f"{{{NS_A}}}r"):
                    t = r.find(f"{{{NS_A}}}t")
                    if t is not None:
                        t.text = ""
                # 也清空段落级别的文本
                for t in p.findall(f"{{{NS_A}}}t"):
                    t.text = ""

    def _convert_to_placeholder(self, shape_el: etree.Element, idx: int) -> None:
        """将形状转换为占位符。"""
        # 找到nvSpPr或nvCxnSpPr
        nv_pr = shape_el.find(_ns("nvSpPr"))
        if nv_pr is None:
            nv_pr = shape_el.find(_ns("nvCxnSpPr"))
        if nv_pr is None:
            return

        # 找到nvPr（如果不存在则创建）
        nvPr = nv_pr.find(_ns("nvPr"))
        if nvPr is None:
            nvPr = etree.SubElement(nv_pr, _ns("nvPr"))

        # 检查是否已经有ph
        existing_ph = nvPr.find(_ns("ph"))
        if existing_ph is not None:
            return

        # 添加ph元素到nvPr (使用p命名空间，与原始layout一致)
        ph = etree.SubElement(nvPr, _ns("ph"))
        # 第一个形状设为标题，其他设为body
        if idx == 0:
            ph.set("type", "title")
        else:
            ph.set("idx", str(idx))
            ph.set("type", "body")

    def _update_content_types(self, extract_dir: Path, layout_num: int) -> None:
        """更新[Content_Types].xml，添加新的slideLayout。"""
        ct_path = extract_dir / "[Content_Types].xml"
        ct_root = etree.parse(str(ct_path)).getroot()

        # 检查是否已存在
        part_name = f"/ppt/slideLayouts/slideLayout{layout_num}.xml"
        existing = ct_root.find(f".//{{{NS_CT}}}Override[@PartName='{part_name}']")
        if existing is not None:
            return

        override = etree.SubElement(ct_root, _ns("Override", NS_CT))
        override.set("PartName", part_name)
        override.set(
            "ContentType",
            "application/vnd.openxmlformats-officedocument.presentationml.slideLayout+xml"
        )

        etree.ElementTree(ct_root).write(
            str(ct_path), xml_declaration=True, encoding="UTF-8", standalone=True
        )

    def _update_slide_master(self, extract_dir: Path, layout_num: int) -> int:
        """更新slideMaster的sldLayoutIdLst，添加新layout引用。"""
        master_xml_path = extract_dir / "ppt" / "slideMasters" / "slideMaster1.xml"
        master_root = etree.parse(str(master_xml_path)).getroot()

        sldLayoutIdLst = master_root.find(_ns("sldLayoutIdLst"))

        # 生成新的ID（最大ID + 1）
        max_id = 0
        for layout_id in sldLayoutIdLst:
            lid = int(layout_id.get("id", 0))
            if lid > max_id:
                max_id = lid
        new_id = max_id + 1

        # 添加新的sldLayoutId
        new_layout_id = etree.SubElement(sldLayoutIdLst, _ns("sldLayoutId"))
        new_layout_id.set("id", str(new_id))

        # 需要添加关系引用
        master_rels_path = extract_dir / "ppt" / "slideMasters" / "_rels" / "slideMaster1.xml.rels"
        rels_root = etree.parse(str(master_rels_path)).getroot()

        # 找到最大的rId
        max_rid = 0
        for rel in rels_root:
            rid = rel.get("Id", "")
            if rid.startswith("rId"):
                try:
                    num = int(rid[3:])
                    if num > max_rid:
                        max_rid = num
                except ValueError:
                    pass

        new_rid = f"rId{max_rid + 1}"
        new_layout_id.set(f"{{{NS_R}}}id", new_rid)

        # 添加关系
        new_rel = etree.SubElement(rels_root, _ns("Relationship", NS_RELS))
        new_rel.set("Id", new_rid)
        new_rel.set(
            "Type",
            "http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideLayout"
        )
        new_rel.set("Target", f"../slideLayouts/slideLayout{layout_num}.xml")

        # 保存
        etree.ElementTree(master_root).write(
            str(master_xml_path), xml_declaration=True, encoding="UTF-8", standalone=True
        )
        etree.ElementTree(rels_root).write(
            str(master_rels_path), xml_declaration=True, encoding="UTF-8", standalone=True
        )

        return len(sldLayoutIdLst) - 1  # 返回新layout的索引

    def _repack_pptx(self, extract_dir: Path, output_path: Path) -> None:
        """重新打包为pptx文件。"""
        output_path.parent.mkdir(parents=True, exist_ok=True)
        with zipfile.ZipFile(output_path, 'w', zipfile.ZIP_DEFLATED) as zf:
            for file_path in extract_dir.rglob("*"):
                if file_path.is_file():
                    arcname = file_path.relative_to(extract_dir)
                    zf.write(file_path, arcname)


# ========== CLI 入口 ==========

def main(argv: Optional[List[str]] = None) -> int:
    parser = argparse.ArgumentParser(description="Layout导入器")
    parser.add_argument("--source", required=True, type=Path, help="源PPT路径")
    parser.add_argument("--target", type=Path, help="目标模板路径")
    parser.add_argument("--analyze", action="store_true", help="仅分析源模板")
    parser.add_argument("--slides", type=str, help="要导入的slide索引（逗号分隔）")
    parser.add_argument("--auto", action="store_true", help="自动导入最有代表性的slides")
    parser.add_argument("--top", type=int, default=5, help="自动导入时选择前N个")
    parser.add_argument("--verbose", action="store_true", help="详细日志")
    args = parser.parse_args(argv)

    if args.verbose:
        logging.getLogger().setLevel(logging.DEBUG)

    try:
        importer = LayoutImporter(args.source, args.target or args.source)

        if args.analyze:
            candidates = importer.analyze_source()
            print(json.dumps(candidates[:10], ensure_ascii=False, indent=2))
            return 0

        if not args.target:
            print("错误: --target 是导入必需的参数")
            return 1

        if args.auto:
            candidates = importer.analyze_source()
            slide_indices = [c["index"] for c in candidates[:args.top]]
            print(f"自动选择 {len(slide_indices)} 个slides: {slide_indices}")
        elif args.slides:
            slide_indices = [int(x.strip()) for x in args.slides.split(",")]
        else:
            print("错误: 请指定 --slides 或 --auto")
            return 1

        result = importer.import_slides(slide_indices)
        print(json.dumps(result, ensure_ascii=False, indent=2))
        return 0

    except Exception as exc:
        logger.error("导入失败: %s", exc)
        import traceback
        traceback.print_exc()
        print(f"❌ 导入失败: {exc}")
        return 1


if __name__ == "__main__":
    sys.exit(main())
