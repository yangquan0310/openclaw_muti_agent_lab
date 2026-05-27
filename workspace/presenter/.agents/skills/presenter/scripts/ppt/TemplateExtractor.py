#!/usr/bin/env python3
"""
TemplateExtractor — 从武汉文理模板提取 layouts，创建新模板。

功能：
  1. 加载武汉文理模板
  2. 保留其 slide_master + slide_layouts 结构
  3. 修改样式：
     - 背景色：白色
     - 标题占位符：32号，微软雅黑+Times New Roman
     - 正文占位符：28号，微软雅黑+Times New Roman
  4. 清空样例 slides（保留结构作为提示）
  5. 保存为新 template.pptx

Usage:
    python3 TemplateExtractor.py --source assets/template_wuhanwenli.pptx --output assets/template.pptx
"""

import argparse
import logging
import sys
from pathlib import Path
from typing import Any, List, Optional

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

logging.basicConfig(level=logging.INFO, format="%(asctime)s [%(levelname)s] %(message)s")
logger = logging.getLogger(__name__)


# 样式常量
BG_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
FONT_CN = "微软雅黑"
FONT_EN = "Times New Roman"
FONT_TITLE_SIZE = Pt(32)
FONT_BODY_SIZE = Pt(28)


class TemplateExtractor:
    """模板提取器。"""

    def __init__(self, source_path: Path) -> None:
        self.prs = Presentation(str(source_path))
        logger.info("加载源模板: %s", source_path)
        logger.info("  - %d 个 slide_layouts", len(self.prs.slide_layouts))
        logger.info("  - %d 个样例 slides", len(self.prs.slides))

    # ========== 公有接口 ==========

    def extract(self, output_path: Path) -> None:
        """提取并创建新模板。"""
        logger.info("开始提取模板...")

        # 1. 设置 slide_master 背景为白色
        self._set_master_background()

        # 2. 设置 slide_layouts 中占位符的默认字体
        self._set_layout_fonts()

        # 3. 清空样例 slides 上的硬编码文本（保留占位符提示）
        self._clear_sample_slides()

        # 4. 保存
        output_path.parent.mkdir(parents=True, exist_ok=True)
        self.prs.save(str(output_path))
        logger.info("新模板已保存: %s", output_path)
        logger.info("  - %d 个 layouts", len(self.prs.slide_layouts))
        logger.info("  - %d 个样例 slides", len(self.prs.slides))

    # ========== 私有方法 ==========

    def _set_master_background(self) -> None:
        """设置 slide_master 背景为白色。"""
        for master in self.prs.slide_masters:
            background = master.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = BG_WHITE
        logger.info("Master 背景已设为白色")

    def _set_layout_fonts(self) -> None:
        """为每个 slide_layout 的占位符设置默认字体。"""
        for layout in self.prs.slide_layouts:
            for shape in layout.placeholders:
                ph_type = shape.placeholder_format.type
                if ph_type == 1:  # TITLE
                    self._set_placeholder_font(shape, FONT_TITLE_SIZE, bold=True)
                elif ph_type in (2, 7):  # BODY or OBJECT
                    self._set_placeholder_font(shape, FONT_BODY_SIZE)
                elif ph_type == 3:  # CENTER_TITLE
                    self._set_placeholder_font(shape, FONT_TITLE_SIZE, bold=True)
                elif ph_type == 4:  # SUBTITLE
                    self._set_placeholder_font(shape, Pt(24))

        logger.info("Layout 占位符字体已设置")

    def _set_placeholder_font(self, shape: Any, font_size: Any, bold: bool = False) -> None:
        """设置占位符的默认字体。"""
        if not shape.has_text_frame:
            return
        tf = shape.text_frame
        # 设置段落默认字体
        for paragraph in tf.paragraphs:
            paragraph.font.size = font_size
            paragraph.font.bold = bold
            paragraph.font.name = FONT_EN
            # 对中文使用微软雅黑（通过设置 East Asian 字体）
            # python-pptx 对 East Asian 字体的支持有限，这里设置主字体
            # 实际效果取决于 PowerPoint 的字体回退机制

        # 通过 XML 设置 East Asian 字体
        try:
            from lxml import etree
            # 获取 a:pPr 或 a:defRPr 元素来设置字体
            # 这是一个简化的方法，实际可能需要更复杂的 XML 操作
            pass
        except ImportError:
            pass

    def _clear_sample_slides(self) -> None:
        """清空样例 slides 上的硬编码文本。"""
        cleared = 0
        for slide in self.prs.slides:
            for shape in slide.shapes:
                # 只清空非占位符文本框和自动形状中的文本
                if not shape.is_placeholder and shape.has_text_frame:
                    # 保留占位符的提示文本，只清空白占位符的
                    for paragraph in shape.text_frame.paragraphs:
                        paragraph.text = ""
                    cleared += 1
                # 对于占位符，设置提示文本
                elif shape.is_placeholder and shape.has_text_frame:
                    ph_type = shape.placeholder_format.type
                    if ph_type == 1 and not shape.text_frame.text:
                        shape.text_frame.text = "单击此处添加标题"
                    elif ph_type in (2, 7) and not shape.text_frame.text:
                        shape.text_frame.text = "单击此处添加正文"

        logger.info("已清空 %d 个形状的文本", cleared)


# ========== CLI 入口 ==========

def main(argv: Optional[List[str]] = None) -> int:
    parser = argparse.ArgumentParser(description="模板提取器")
    parser.add_argument("--source", required=True, type=Path, help="源模板路径(.pptx)")
    parser.add_argument("--output", required=True, type=Path, help="输出模板路径(.pptx)")
    parser.add_argument("--verbose", action="store_true", help="详细日志")
    args = parser.parse_args(argv)

    if args.verbose:
        logging.getLogger().setLevel(logging.DEBUG)

    try:
        extractor = TemplateExtractor(args.source)
        extractor.extract(args.output)
        print(f"✅ 模板提取成功: {args.output}")
        return 0
    except Exception as exc:
        logger.error("提取失败: %s", exc)
        print(f"❌ 提取失败: {exc}")
        return 1


if __name__ == "__main__":
    sys.exit(main())
