#!/usr/bin/env python3
"""
PPTX 编译器 — Layout 是画布，Slide 设计是核心。

核心原则：
  1. Layout 提供背景、字体、颜色、页脚 → 直接使用
  2. Slide 的内容呈现 → 编译器根据 @structure 自主设计排版
  3. Placeholder 只作为位置参考，不强制填充
  4. 大部分内容通过 add_textbox/add_shape 自主放置

Usage:
    python3 PptxCompiler.py --template template.pptx --script script.json --output out.pptx
"""

import argparse
import json
import logging
import re
import sys
from pathlib import Path
from typing import Any, Dict, List, Optional

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

logging.basicConfig(level=logging.INFO, format="%(asctime)s [%(levelname)s] %(message)s")
logger = logging.getLogger(__name__)


# ========== 配色方案 ==========

# 多套主题配色
THEMES = {
    "blue": {
        "primary": RGBColor(0x00, 0x96, 0xC7),
        "secondary": RGBColor(0x48, 0xCA, 0xE4),
        "accent": RGBColor(0x90, 0xE0, 0xEF),
        "dark": RGBColor(0x21, 0x25, 0x29),
        "light": RGBColor(0xF8, 0xF9, 0xFA),
        "white": RGBColor(0xFF, 0xFF, 0xFF),
        "gray": RGBColor(0xCC, 0xCC, 0xCC),
        "orange": RGBColor(0xF4, 0xA2, 0x61),
        "purple": RGBColor(0x9B, 0x59, 0xB6),
        "bg_dark": RGBColor(0x0D, 0x1B, 0x2A),
        "bg_card": RGBColor(0x1B, 0x2A, 0x3D),
    },
    "green": {
        "primary": RGBColor(0x27, 0xAE, 0x60),
        "secondary": RGBColor(0x2E, 0xCC, 0x71),
        "accent": RGBColor(0x82, 0xE0, 0xAA),
        "dark": RGBColor(0x1E, 0x3A, 0x2A),
        "light": RGBColor(0xF4, 0xF9, 0xF5),
        "white": RGBColor(0xFF, 0xFF, 0xFF),
        "gray": RGBColor(0x95, 0xA5, 0xA6),
        "orange": RGBColor(0xE6, 0x7E, 0x22),
        "purple": RGBColor(0x8E, 0x44, 0xAD),
        "bg_dark": RGBColor(0x1A, 0x2E, 0x1C),
        "bg_card": RGBColor(0x2D, 0x4A, 0x34),
    },
    "purple": {
        "primary": RGBColor(0x9B, 0x59, 0xB6),
        "secondary": RGBColor(0xBB, 0x8F, 0xCE),
        "accent": RGBColor(0xD7, 0xBD, 0xE8),
        "dark": RGBColor(0x2C, 0x1E, 0x3A),
        "light": RGBColor(0xF9, 0xF5, 0xFC),
        "white": RGBColor(0xFF, 0xFF, 0xFF),
        "gray": RGBColor(0xAA, 0x9D, 0xB5),
        "orange": RGBColor(0xF4, 0xA2, 0x61),
        "purple": RGBColor(0x9B, 0x59, 0xB6),
        "bg_dark": RGBColor(0x1A, 0x12, 0x2A),
        "bg_card": RGBColor(0x2D, 0x23, 0x4A),
    },
    "orange": {
        "primary": RGBColor(0xE6, 0x7E, 0x22),
        "secondary": RGBColor(0xF5, 0xA6, 0x23),
        "accent": RGBColor(0xFC, 0xE4, 0xB5),
        "dark": RGBColor(0x3D, 0x2B, 0x1A),
        "light": RGBColor(0xFD, 0xF6, 0xED),
        "white": RGBColor(0xFF, 0xFF, 0xFF),
        "gray": RGBColor(0xAA, 0x95, 0x7D),
        "orange": RGBColor(0xE6, 0x7E, 0x22),
        "purple": RGBColor(0x9B, 0x59, 0xB6),
        "bg_dark": RGBColor(0x2A, 0x1E, 0x12),
        "bg_card": RGBColor(0x4A, 0x34, 0x23),
    },
    "gray": {
        "primary": RGBColor(0x5D, 0x6D, 0x7E),
        "secondary": RGBColor(0x85, 0x9B, 0xB5),
        "accent": RGBColor(0xBD, 0xC3, 0xC7),
        "dark": RGBColor(0x2C, 0x3E, 0x50),
        "light": RGBColor(0xF4, 0xF6, 0xF7),
        "white": RGBColor(0xFF, 0xFF, 0xFF),
        "gray": RGBColor(0x95, 0xA5, 0xA6),
        "orange": RGBColor(0xE6, 0x7E, 0x22),
        "purple": RGBColor(0x9B, 0x59, 0xB6),
        "bg_dark": RGBColor(0x2C, 0x3E, 0x50),
        "bg_card": RGBColor(0x34, 0x49, 0x5E),
    },
}


class Theme:
    """配色方案类，支持多主题切换。"""

    def __init__(self, theme_name: str = "blue") -> None:
        if theme_name not in THEMES:
            logger.warning("未知主题 '%s'，使用 'blue' 默认", theme_name)
            theme_name = "blue"
        t = THEMES[theme_name]
        self.PRIMARY = t["primary"]
        self.SECONDARY = t["secondary"]
        self.ACCENT = t["accent"]
        self.DARK = t["dark"]
        self.LIGHT = t["light"]
        self.WHITE = t["white"]
        self.GRAY = t["gray"]
        self.ORANGE = t["orange"]
        self.PURPLE = t["purple"]
        self.BG_DARK = t["bg_dark"]
        self.BG_CARD = t["bg_card"]


class PptxCompiler:
    """PPTX 编译器：Layout 是画布，Slide 设计是核心。"""

    def __init__(self, template_path: Path, theme: str = "blue") -> None:
        self.prs = Presentation(str(template_path))
        self._page_num = 0
        self._layouts: List[Any] = list(self.prs.slide_layouts)
        self._theme = Theme(theme)
        logger.info("模板加载: %d 个 slide_layouts, 主题: %s", len(self._layouts), theme)

    # ========== 公有接口 ==========

    def compile_page(self, layout_idx: int, lines: List[str],
                     structure: str = "list", chart_design: str = "") -> None:
        """编译单页：获取 layout 画布，自主设计 slide。"""
        if layout_idx >= len(self._layouts):
            logger.warning("索引 %d 越界，回退到 layout 0", layout_idx)
            layout_idx = 0

        layout = self._layouts[layout_idx]
        slide = self.prs.slides.add_slide(layout)
        self._page_num += 1

        # 删除标题和内容占位符，避免干扰自主设计
        self._remove_content_placeholders(slide)

        # 根据 structure 自主设计
        if structure == "cover":
            self._design_cover(slide, lines)
        elif structure == "section":
            self._design_section(slide, lines)
        elif structure == "compare":
            self._design_compare(slide, lines)
        elif structure == "timeline":
            self._design_timeline(slide, lines)
        elif structure == "flowchart":
            self._design_flowchart(slide, lines)
        elif structure == "cards":
            self._design_cards(slide, lines)
        elif structure == "table":
            self._design_table(slide, lines)
        elif structure == "toc":
            self._design_toc(slide, lines)
        elif structure == "quote":
            self._design_quote(slide, lines)
        elif structure == "stats":
            self._design_stats(slide, lines)
        else:
            # 默认：标准列表布局
            self._design_list(slide, lines)

    def save(self, output_path: Path) -> None:
        output_path.parent.mkdir(parents=True, exist_ok=True)
        self.prs.save(str(output_path))
        logger.info("PPT 已保存: %s (共 %d 页)", output_path, len(self.prs.slides))

    def list_layouts(self) -> Dict[str, List[Dict[str, Any]]]:
        layouts_info = []
        for idx, layout in enumerate(self._layouts):
            placeholders = []
            for shape in layout.placeholders:
                text = shape.text_frame.text.strip()[:50] if hasattr(shape, "text_frame") else ""
                placeholders.append({
                    "idx": shape.placeholder_format.idx,
                    "type": str(shape.placeholder_format.type),
                    "name": shape.name,
                    "text_preview": text,
                })
            layouts_info.append({
                "index": idx,
                "name": layout.name,
                "placeholders_count": len(placeholders),
                "placeholders": placeholders,
            })
        return {"layouts": layouts_info}

    # ========== 清空占位符 ==========

    def _remove_content_placeholders(self, slide: Any) -> None:
        """删除标题和内容占位符，避免干扰自主设计。

        保留日期(10)/页脚(11)/编号(12)占位符，它们显示母版页脚信息。
        """
        for shape in list(slide.shapes):
            if not shape.is_placeholder:
                continue
            ph_type = shape.placeholder_format.type
            # 删除 title(1)/center_title(3)/body(2)/subtitle(4)/obj(7)
            # 保留 date(16)/footer(15)/slide_num(13)
            if ph_type in (1, 2, 3, 4, 7):
                sp = shape.element
                sp.getparent().remove(sp)

    # ========== 封面设计 ==========

    def _design_cover(self, slide: Any, lines: List[str]) -> None:
        """封面设计：大标题居中，信息层级分明。"""
        if not lines:
            return

        # 解析内容
        title = lines[0] if len(lines) > 0 else ""
        subtitle = lines[1] if len(lines) > 1 else ""
        teacher = lines[2] if len(lines) > 2 else ""
        institution = lines[3] if len(lines) > 3 else ""

        # 课程名（最大）
        if title:
            self._add_text_box(slide, Inches(0), Inches(2.5), Inches(10), Inches(1.2),
                               title, "微软雅黑", 44, font_color=self._theme.PRIMARY,
                               alignment=PP_ALIGN.CENTER, bold=True)

        # 章节名（次大）
        if subtitle:
            self._add_text_box(slide, Inches(0), Inches(3.8), Inches(10), Inches(0.8),
                               subtitle, "微软雅黑", 32, font_color=self._theme.DARK,
                               alignment=PP_ALIGN.CENTER, bold=True)

        # 教师信息
        if teacher:
            self._add_text_box(slide, Inches(0), Inches(5.0), Inches(10), Inches(0.5),
                               teacher, "微软雅黑", 20, font_color=self._theme.DARK,
                               alignment=PP_ALIGN.CENTER)

        # 机构信息
        if institution:
            self._add_text_box(slide, Inches(0), Inches(5.6), Inches(10), Inches(0.4),
                               institution, "微软雅黑", 16, font_color=self._theme.GRAY,
                               alignment=PP_ALIGN.CENTER)

        # 装饰线
        if title:
            line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                          Inches(3.5), Inches(3.6), Inches(3.0), Inches(0.04))
            line.fill.solid()
            line.fill.fore_color.rgb = self._theme.PRIMARY
            line.line.fill.background()

    # ========== 章节分隔设计 ==========

    def _design_section(self, slide: Any, lines: List[str]) -> None:
        """章节分隔：大字+简洁装饰，视觉区分。"""
        if not lines:
            return

        part = lines[0] if len(lines) > 0 else ""
        title = lines[1] if len(lines) > 1 else ""
        keywords = lines[2] if len(lines) > 2 else ""

        # PART 编号（大字）
        if part:
            self._add_text_box(slide, Inches(1.0), Inches(2.0), Inches(8.0), Inches(1.0),
                               part, "微软雅黑", 48, font_color=self._theme.PRIMARY,
                               alignment=PP_ALIGN.LEFT, bold=True)

        # 章节标题
        if title:
            self._add_text_box(slide, Inches(1.0), Inches(3.2), Inches(8.0), Inches(0.8),
                               title, "微软雅黑", 36, font_color=self._theme.DARK,
                               alignment=PP_ALIGN.LEFT, bold=True)

        # 关键词
        if keywords:
            self._add_text_box(slide, Inches(1.0), Inches(4.2), Inches(8.0), Inches(0.5),
                               keywords, "微软雅黑", 20, font_color=self._theme.GRAY,
                               alignment=PP_ALIGN.LEFT)

        # 左侧竖线装饰
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                     Inches(0.5), Inches(2.0), Inches(0.08), Inches(2.5))
        bar.fill.solid()
        bar.fill.fore_color.rgb = self._theme.PRIMARY
        bar.line.fill.background()

    # ========== 标准列表设计 ==========

    def _design_list(self, slide: Any, lines: List[str]) -> None:
        """标准列表设计：标题 + bullet 要点。"""
        if not lines:
            return

        title = lines[0]
        items = [re.sub(r"^[-*•\d\.\s]+", "", line.strip())
                 for line in lines[1:] if line.strip()]

        # 标题
        self._add_text_box(slide, Inches(0.8), Inches(0.6), Inches(8.4), Inches(0.7),
                           title, "微软雅黑", 32, font_color=self._theme.DARK, bold=True)

        # 分隔线
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                      Inches(0.8), Inches(1.4), Inches(1.5), Inches(0.04))
        line.fill.solid()
        line.fill.fore_color.rgb = self._theme.PRIMARY
        line.line.fill.background()

        # 要点列表
        y = Inches(1.8)
        for i, item in enumerate(items):
            # 编号圆点
            dot = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                         Inches(0.9), y + Inches(0.08), Inches(0.2), Inches(0.2))
            dot.fill.solid()
            dot.fill.fore_color.rgb = self._theme.PRIMARY
            dot.line.fill.background()

            # 编号文字
            self._add_text_box(slide, Inches(0.85), y + Inches(0.05), Inches(0.3), Inches(0.3),
                               str(i + 1), "微软雅黑", 12, font_color=self._theme.WHITE,
                               alignment=PP_ALIGN.CENTER)

            # 内容
            self._add_text_box(slide, Inches(1.3), y, Inches(7.5), Inches(0.5),
                               item, "微软雅黑", 22, font_color=self._theme.DARK)
            y += Inches(0.7)

    # ========== 双栏对比设计 ==========

    def _design_compare(self, slide: Any, lines: List[str]) -> None:
        """双栏对比设计：左右对称，中间 vs 标识。"""
        if not lines:
            return

        # 解析左右两侧内容
        left_label = ""
        left_items = []
        right_label = ""
        right_items = []
        current_side = None

        for line in lines:
            stripped = line.strip()
            if not stripped:
                continue

            # 检测左侧标记
            if "左侧" in stripped or "左边" in stripped:
                current_side = "left"
                # 提取标签（去掉"左侧""左边"前缀）
                label = re.sub(r"^.*?[：:]\s*", "", stripped)
                if label:
                    left_label = label
                continue

            # 检测右侧标记
            if "右侧" in stripped or "右边" in stripped:
                current_side = "right"
                label = re.sub(r"^.*?[：:]\s*", "", stripped)
                if label:
                    right_label = label
                continue

            # 当前侧的内容
            cleaned = re.sub(r"^[-*•\d\.\s]+", "", stripped)
            if cleaned:
                if current_side == "left":
                    left_items.append(cleaned)
                elif current_side == "right":
                    right_items.append(cleaned)

        # 标题（用左右标签组合）
        title = f"{left_label}  vs  {right_label}" if left_label and right_label else ""
        if title:
            self._add_text_box(slide, Inches(0.8), Inches(0.6), Inches(8.4), Inches(0.7),
                               title, "微软雅黑", 32, font_color=self._theme.DARK, bold=True)

        # 左侧卡片
        left_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                           Inches(0.8), Inches(1.5), Inches(4.0), Inches(4.5))
        left_card.fill.solid()
        left_card.fill.fore_color.rgb = self._theme.LIGHT
        left_card.line.color.rgb = self._theme.PRIMARY
        left_card.line.width = Pt(2)

        # 右侧卡片
        right_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                            Inches(5.2), Inches(1.5), Inches(4.0), Inches(4.5))
        right_card.fill.solid()
        right_card.fill.fore_color.rgb = self._theme.LIGHT
        right_card.line.color.rgb = self._theme.ORANGE
        right_card.line.width = Pt(2)

        # 中间 vs
        self._add_text_box(slide, Inches(4.5), Inches(3.5), Inches(1.0), Inches(0.5),
                           "VS", "微软雅黑", 24, font_color=self._theme.GRAY,
                           alignment=PP_ALIGN.CENTER, bold=True)

        # 填充左侧内容
        y = Inches(1.8)
        for item in left_items:
            self._add_text_box(slide, Inches(1.0), y, Inches(3.6), Inches(0.5),
                               item, "微软雅黑", 18, font_color=self._theme.DARK)
            y += Inches(0.6)

        # 填充右侧内容
        y = Inches(1.8)
        for item in right_items:
            self._add_text_box(slide, Inches(5.4), y, Inches(3.6), Inches(0.5),
                               item, "微软雅黑", 18, font_color=self._theme.DARK)
            y += Inches(0.6)

    # ========== 时间轴设计 ==========

    def _design_timeline(self, slide: Any, lines: List[str]) -> None:
        """时间轴设计：标题 + 水平时间轴 + 节点。"""
        if not lines:
            return

        title = lines[0]

        # 提取时间节点 [时间] 描述
        nodes = []
        for line in lines[1:]:
            m = re.match(r"\[(.+?)\]\s*(.+)", line.strip())
            if m:
                nodes.append((m.group(1).strip(), m.group(2).strip()))

        # 标题
        self._add_text_box(slide, Inches(0.8), Inches(0.6), Inches(8.4), Inches(0.7),
                           title, "微软雅黑", 32, font_color=self._theme.DARK, bold=True)

        if len(nodes) < 2:
            # 无时间节点，用标准列表 fallback
            self._design_list(slide, lines)
            return

        # 绘制时间轴
        n = len(nodes)
        start_x = Inches(1.0)
        end_x = Inches(9.5)
        total_w = end_x - start_x
        spacing = total_w / (n + 1)
        y_line = Inches(4.0)
        y_node = Inches(3.7)
        y_time = Inches(3.0)
        y_desc = Inches(4.3)

        colors = [self._theme.PRIMARY, self._theme.SECONDARY, self._theme.ACCENT, self._theme.ORANGE, self._theme.PURPLE]

        # 水平主线
        line_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                            start_x, y_line, total_w, Inches(0.04))
        line_shape.fill.solid()
        line_shape.fill.fore_color.rgb = self._theme.PRIMARY
        line_shape.line.fill.background()

        # 节点
        for i, (time, desc) in enumerate(nodes):
            cx = start_x + spacing * (i + 1)
            color = colors[i % len(colors)]

            # 节点圆
            dot = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                         cx - Inches(0.18), y_node, Inches(0.36), Inches(0.36))
            dot.fill.solid()
            dot.fill.fore_color.rgb = color
            dot.line.color.rgb = self._theme.WHITE
            dot.line.width = Pt(3)

            # 时间标签（上方）
            self._add_text_box(slide, cx - Inches(1.0), y_time, Inches(2.0), Inches(0.5),
                               time, "微软雅黑", 14, font_color=color,
                               alignment=PP_ALIGN.CENTER, bold=True)

            # 描述标签（下方）
            short_desc = desc[:15] + "…" if len(desc) > 15 else desc
            self._add_text_box(slide, cx - Inches(1.0), y_desc, Inches(2.0), Inches(0.6),
                               short_desc, "微软雅黑", 11, font_color=self._theme.DARK,
                               alignment=PP_ALIGN.CENTER)

    # ========== 流程图设计 ==========

    def _design_flowchart(self, slide: Any, lines: List[str]) -> None:
        """流程图设计：标题 + 水平步骤框 + 箭头。"""
        if not lines:
            return

        title = lines[0]

        # 提取步骤
        steps = []
        for line in lines[1:]:
            m = re.match(r"(?:步骤\s*(\d+)[：:]\s*)(.+)", line.strip())
            if m:
                steps.append((m.group(1), m.group(2).strip()))
            else:
                m2 = re.match(r"(\d+)\.\s*(.+)", line.strip())
                if m2:
                    steps.append((m2.group(1), m2.group(2).strip()))

        if not steps:
            # 无步骤格式，fallback
            self._design_list(slide, lines)
            return

        # 标题
        self._add_text_box(slide, Inches(0.8), Inches(0.6), Inches(8.4), Inches(0.7),
                           title, "微软雅黑", 32, font_color=self._theme.DARK, bold=True)

        n = len(steps)
        box_w = min(Inches(2.2), Inches(9.0) / n - Inches(0.3))
        box_h = Inches(1.0)
        total_w = n * box_w + (n - 1) * Inches(0.4)
        start_x = (Inches(10.0) - total_w) / 2
        y_top = Inches(2.5)

        colors = [self._theme.PRIMARY, self._theme.SECONDARY, self._theme.ACCENT, self._theme.ORANGE, self._theme.PURPLE]

        for i, (num, text) in enumerate(steps):
            x = start_x + i * (box_w + Inches(0.4))
            color = colors[i % len(colors)]

            # 步骤框
            box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                         x, y_top, box_w, box_h)
            box.fill.solid()
            box.fill.fore_color.rgb = color
            box.line.fill.background()

            # 步骤文字
            short_text = text[:12] + "…" if len(text) > 12 else text
            tf = box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = f"{num}. {short_text}"
            p.alignment = PP_ALIGN.CENTER
            run = p.runs[0]
            run.font.name = "微软雅黑"
            run.font.size = Pt(13)
            run.font.color.rgb = self._theme.WHITE
            run.font.bold = True

            # 箭头
            if i < n - 1:
                arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                               x + box_w + Inches(0.05), y_top + Inches(0.35),
                                               Inches(0.3), Inches(0.3))
                arrow.fill.solid()
                arrow.fill.fore_color.rgb = self._theme.GRAY
                arrow.line.fill.background()

    # ========== 卡片设计 ==========

    def _design_cards(self, slide: Any, lines: List[str]) -> None:
        """卡片网格设计：标题 + 2×2/2×3 卡片。"""
        if not lines:
            return

        title = lines[0]
        items = [re.sub(r"^[-*•\d\.\s]+", "", line.strip())
                 for line in lines[1:] if line.strip()]

        if not items:
            return

        # 标题
        self._add_text_box(slide, Inches(0.8), Inches(0.6), Inches(8.4), Inches(0.7),
                           title, "微软雅黑", 32, font_color=self._theme.DARK, bold=True)

        n = len(items)
        cols = min(n, 3)
        rows = (n + cols - 1) // cols

        card_w = Inches(2.8)
        card_h = Inches(1.2)
        gap_x = Inches(0.3)
        gap_y = Inches(0.25)
        total_w = cols * card_w + (cols - 1) * gap_x
        start_x = (Inches(10.0) - total_w) / 2
        start_y = Inches(1.8)

        colors = [self._theme.PRIMARY, self._theme.SECONDARY, self._theme.ACCENT, self._theme.ORANGE, self._theme.PURPLE]

        for i, item in enumerate(items):
            row = i // cols
            col = i % cols
            x = start_x + col * (card_w + gap_x)
            y = start_y + row * (card_h + gap_y)
            color = colors[i % len(colors)]

            # 卡片背景
            card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                          x, y, card_w, card_h)
            card.fill.solid()
            card.fill.fore_color.rgb = self._theme.LIGHT
            card.line.color.rgb = color
            card.line.width = Pt(2)

            # 左侧色条
            bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                         x, y, Inches(0.08), card_h)
            bar.fill.solid()
            bar.fill.fore_color.rgb = color
            bar.line.fill.background()

            # 文字
            short = item[:20] + "…" if len(item) > 20 else item
            self._add_text_box(slide, x + Inches(0.2), y + Inches(0.2),
                               card_w - Inches(0.3), Inches(0.8),
                               short, "微软雅黑", 16, font_color=self._theme.DARK)

    # ========== 表格设计 ==========

    def _design_table(self, slide: Any, lines: List[str]) -> None:
        """表格设计：从内容提取 key:value 绘制真正表格。"""
        if not lines:
            return

        title = lines[0]

        # 提取表格数据
        rows = []
        for line in lines[1:]:
            m = re.match(r"(.+?)[：:](.+)", line.strip())
            if m:
                key = re.sub(r"^[-*•\d\.\s]+", "", m.group(1).strip())
                val = re.sub(r"^[-*•\d\.\s]+", "", m.group(2).strip())
                if key and val:
                    rows.append((key, val))

        if not rows:
            # 无法解析，fallback
            self._design_list(slide, lines)
            return

        # 标题
        self._add_text_box(slide, Inches(0.8), Inches(0.6), Inches(8.4), Inches(0.7),
                           title, "微软雅黑", 32, font_color=self._theme.DARK, bold=True)

        # 绘制表格
        num_rows = len(rows) + 1
        num_cols = 2
        x = Inches(1.0)
        y = Inches(1.6)
        w = Inches(8.0)
        h = Inches(0.5 + num_rows * 0.55)

        table = slide.shapes.add_table(num_rows, num_cols, x, y, w, h).table

        # 表头
        table.cell(0, 0).text = "类型"
        table.cell(0, 1).text = "说明"
        for col in range(num_cols):
            cell = table.cell(0, col)
            cell.fill.solid()
            cell.fill.fore_color.rgb = self._theme.PRIMARY
            p = cell.text_frame.paragraphs[0]
            run = p.runs[0]
            run.font.color.rgb = self._theme.WHITE
            run.font.bold = True
            run.font.size = Pt(14)
            run.font.name = "微软雅黑"

        # 数据行
        for i, (key, val) in enumerate(rows):
            r = i + 1
            table.cell(r, 0).text = key
            table.cell(r, 1).text = val
            color = self._theme.LIGHT if r % 2 == 1 else self._theme.WHITE
            for col in range(num_cols):
                cell = table.cell(r, col)
                cell.fill.solid()
                cell.fill.fore_color.rgb = color
                p = cell.text_frame.paragraphs[0]
                run = p.runs[0]
                run.font.size = Pt(12)
                run.font.name = "微软雅黑"
                run.font.color.rgb = self._theme.DARK

    # ========== 目录设计 ==========

    def _design_toc(self, slide: Any, lines: List[str]) -> None:
        """目录设计：左侧标题，右侧章节列表。"""
        if not lines:
            return

        title = lines[0] if len(lines) > 0 else "目录"
        items = [re.sub(r"^[-*\u2022\d\.\s]+", "", line.strip())
                 for line in lines[1:] if line.strip()]

        # 左侧标题
        self._add_text_box(slide, Inches(0.8), Inches(0.6), Inches(3.0), Inches(0.7),
                           title, "微软雅黑", 32, font_color=self._theme.PRIMARY, bold=True)

        # 左侧装饰线
        line_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                      Inches(0.8), Inches(1.4), Inches(1.5), Inches(0.04))
        line_shape.fill.solid()
        line_shape.fill.fore_color.rgb = self._theme.PRIMARY
        line_shape.line.fill.background()

        # 右侧章节列表
        y = Inches(1.8)
        for i, item in enumerate(items):
            num_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                             Inches(4.5), y, Inches(0.5), Inches(0.4))
            num_box.fill.solid()
            num_box.fill.fore_color.rgb = self._theme.PRIMARY
            num_box.line.fill.background()

            self._add_text_box(slide, Inches(4.55), y + Inches(0.02), Inches(0.4), Inches(0.4),
                               str(i + 1), "微软雅黑", 14, font_color=self._theme.WHITE,
                               alignment=PP_ALIGN.CENTER)

            self._add_text_box(slide, Inches(5.2), y, Inches(5.0), Inches(0.4),
                               item, "微软雅黑", 18, font_color=self._theme.DARK)

            dot_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                              Inches(5.1), y + Inches(0.18), Inches(0.1), Inches(0.04))
            dot_line.fill.solid()
            dot_line.fill.fore_color.rgb = self._theme.GRAY
            dot_line.line.fill.background()

            y += Inches(0.8)

    # ========== 引用设计 ==========

    def _design_quote(self, slide: Any, lines: List[str]) -> None:
        """引用/金句设计：大字居中，装饰引号。"""
        if not lines:
            return

        quote = lines[0]
        author = lines[1] if len(lines) > 1 else ""

        # 装饰引号（左）
        self._add_text_box(slide, Inches(1.0), Inches(1.5), Inches(1.0), Inches(1.0),
                           "\u201c", "微软雅黑", 72, font_color=self._theme.PRIMARY)

        # 装饰引号（右）
        self._add_text_box(slide, Inches(8.5), Inches(4.0), Inches(1.0), Inches(1.0),
                           "\u201d", "微软雅黑", 72, font_color=self._theme.PRIMARY)

        # 引文内容
        self._add_text_box(slide, Inches(1.5), Inches(2.5), Inches(8.0), Inches(2.0),
                           quote, "微软雅黑", 28, font_color=self._theme.DARK,
                           alignment=PP_ALIGN.CENTER, bold=True)

        # 作者
        if author:
            self._add_text_box(slide, Inches(1.5), Inches(5.0), Inches(8.0), Inches(0.5),
                               "— " + author, "微软雅黑", 18, font_color=self._theme.GRAY,
                               alignment=PP_ALIGN.RIGHT)

        # 底部装饰线
        line_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                      Inches(5.5), Inches(5.5), Inches(2.0), Inches(0.04))
        line_shape.fill.solid()
        line_shape.fill.fore_color.rgb = self._theme.PRIMARY
        line_shape.line.fill.background()

    # ========== 统计设计 ==========

    def _design_stats(self, slide: Any, lines: List[str]) -> None:
        """统计设计：大数字突出显示。"""
        if not lines:
            return

        title = lines[0]

        # 标题
        self._add_text_box(slide, Inches(0.8), Inches(0.6), Inches(8.4), Inches(0.7),
                           title, "微软雅黑", 28, font_color=self._theme.DARK, bold=True)

        # 解析统计项
        stats = []
        for line in lines[1:]:
            line = line.strip()
            if not line:
                continue
            m = re.match(r"^([\d\.]+)[\s\n]+(.+)", line, re.DOTALL)
            if m:
                num = m.group(1).strip()
                desc = m.group(2).strip()
                stats.append((num, desc))

        if not stats:
            self._design_list(slide, lines)
            return

        num_stats = len(stats)
        cols = min(num_stats, 3)
        rows = (num_stats + cols - 1) // cols

        card_w = Inches(10.0 / cols)
        card_h = Inches(4.0 / rows)
        start_x = Inches(1.5)
        start_y = Inches(1.8)

        colors = [self._theme.PRIMARY, self._theme.SECONDARY,
                  self._theme.ACCENT, self._theme.ORANGE,
                  self._theme.PURPLE, self._theme.GRAY]

        for i, (num, desc) in enumerate(stats):
            col = i % cols
            row = i // cols

            x = start_x + col * card_w
            y = start_y + row * card_h

            card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                          x, y, card_w - Inches(0.2), card_h - Inches(0.2))
            card.fill.solid()
            card.fill.fore_color.rgb = self._theme.BG_CARD
            card.line.color.rgb = colors[i % len(colors)]
            card.line.width = Pt(2)

            self._add_text_box(slide, x + Inches(0.2), y + Inches(0.2),
                               card_w - Inches(0.4), card_h - Inches(1.2),
                               num, "Arial", 48, font_color=colors[i % len(colors)], bold=True,
                               alignment=PP_ALIGN.CENTER)

            self._add_text_box(slide, x + Inches(0.2), y + card_h - Inches(1.0),
                               card_w - Inches(0.4), Inches(0.8),
                               desc, "微软雅黑", 14, font_color=self._theme.WHITE,
                               alignment=PP_ALIGN.CENTER)

    # ========== 通用工具 ==========

    def _add_text_box(self, slide: Any, left: Any, top: Any, width: Any, height: Any,
                      text: str, font_name: str, font_size_pt: int,
                      font_color: Optional[RGBColor] = None,
                      bold: bool = False, alignment: int = PP_ALIGN.LEFT) -> Any:
        """添加文本框。"""
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.alignment = alignment
        run = p.runs[0]
        run.font.name = font_name
        run.font.size = Pt(font_size_pt)
        run.font.bold = bold
        if font_color:
            run.font.color.rgb = font_color
        return txBox


# ========== CLI 入口 ==========

def _parse_args(argv: Optional[List[str]] = None) -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="PPTX 编译器（Layout 是画布，Slide 设计是核心）")
    parser.add_argument("--template", required=True, type=Path, help="模板PPT路径")
    parser.add_argument("--script", type=Path, help="脚本 JSON 路径")
    parser.add_argument("--output", type=Path, help="输出 PPT 路径")
    parser.add_argument("--list-layouts", action="store_true", help="列出模板 layouts")
    parser.add_argument("--verbose", action="store_true", help="详细日志")
    return parser.parse_args(argv)


def main(argv: Optional[List[str]] = None) -> int:
    args = _parse_args(argv)
    if args.verbose:
        logging.getLogger().setLevel(logging.DEBUG)
    try:
        compiler = PptxCompiler(args.template)

        if args.list_layouts:
            info = compiler.list_layouts()
            print(json.dumps(info, ensure_ascii=False, indent=2))
            return 0

        if not args.script or not args.output:
            print("错误: --script 和 --output 是必需的参数")
            return 1

        data = json.loads(args.script.read_text(encoding="utf-8"))
        pages = data.get("pages", [])

        for page in pages:
            compiler.compile_page(
                int(page.get("layout", 0)),
                page.get("lines", []),
                structure=page.get("structure", "list"),
                chart_design=page.get("chart_design", "")
            )

        compiler.save(args.output)
        print(json.dumps({
            "success": True, "output": str(args.output), "pages": len(pages)
        }, ensure_ascii=False, indent=2))
        return 0
    except Exception as exc:
        logger.error("编译失败: %s", exc)
        import traceback
        traceback.print_exc()
        print(json.dumps({"success": False, "error": str(exc)}, ensure_ascii=False))
        return 1


if __name__ == "__main__":
    sys.exit(main())
