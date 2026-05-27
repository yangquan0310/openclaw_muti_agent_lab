#!/usr/bin/env python3
"""
TemplateBuilder — 创建真正意义的课件模板PPT。

核心设计：
  1. slide_masters: 定义全局背景色、字体、主题
  2. slide_layouts: 每种布局包含正确的 placeholders + 装饰形状
  3. 样例 slides: 占位符显示提示文本（非硬编码内容）

在 slide_layout 层面添加的装饰形状，会自动出现在所有基于该 layout 的新 slide 上。

Usage:
    python3 TemplateBuilder.py --output assets/template_new.pptx
"""

import argparse
import logging
import sys
from pathlib import Path
from typing import Any, List, Optional

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

logging.basicConfig(level=logging.INFO, format="%(asctime)s [%(levelname)s] %(message)s")
logger = logging.getLogger(__name__)


class TemplateBuilder:
    """课件模板构建器。"""

    # 配色系统
    BG_DARK = RGBColor(0x0D, 0x1B, 0x2A)       # 深蓝背景 #0D1B2A
    BG_CARD = RGBColor(0x1B, 0x2A, 0x3D)       # 卡片背景 #1B2A3D
    ACCENT = RGBColor(0x00, 0x96, 0xC7)        # 强调色 #0096C7
    ACCENT_LIGHT = RGBColor(0x48, 0xCA, 0xE4)  # 浅强调 #48CAE4
    TEXT_WHITE = RGBColor(0xFF, 0xFF, 0xFF)    # 纯白
    TEXT_GRAY = RGBColor(0xB0, 0xBE, 0xC5)     # 灰色提示
    ORANGE = RGBColor(0xE6, 0x73, 0x00)        # 橙色点缀

    def __init__(self) -> None:
        self.prs: Any = Presentation()
        self._set_slide_size()

    # ========== 公有接口 ==========

    def build(self, output_path: Path) -> None:
        """构建完整模板并保存。"""
        logger.info("开始构建模板...")

        # 1. 配置 slide_master
        self._config_master()

        # 2. 创建样例 slides（显示提示文本）
        self._create_sample_slides()

        # 4. 保存
        output_path.parent.mkdir(parents=True, exist_ok=True)
        self.prs.save(str(output_path))
        logger.info("模板已保存: %s", output_path)
        logger.info("  - %d 个 slide_layouts", len(self.prs.slide_layouts))
        logger.info("  - %d 个样例 slides", len(self.prs.slides))

    # ========== 私有方法 ==========

    def _set_slide_size(self) -> None:
        """设置 16:9 宽屏。"""
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)
        logger.info("幻灯片尺寸: 16:9 (%.1f x %.1f 英寸)",
                    self.prs.slide_width.inches, self.prs.slide_height.inches)

    def _config_master(self) -> None:
        """配置 slide_master：背景、字体。"""
        master = self.prs.slide_masters[0]

        # 设置背景为深蓝
        background = master.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.BG_DARK
        logger.info("Master 背景已设为深蓝 #0D1B2A")

        # 设置默认字体（通过主题）
        # python-pptx 对主题字体的支持有限，这里仅设置背景

    def _create_sample_slides(self) -> None:
        """创建样例 slides，placeholder 显示提示文本。"""
        layouts = list(self.prs.slide_layouts)

        # 为每个关键 layout 创建一个样例 slide
        sample_configs = [
            (0, "标题幻灯片", "封面页：课程名称、章节、教师信息"),
            (1, "标题和内容", "正文页：标题 + 内容列表"),
            (2, "节标题", "章节分隔页：PART 01 / 章节名称"),
            (3, "两栏内容", "双栏页：左右对比内容"),
            (6, "空白", "空白页：完全自定义内容"),
        ]

        for layout_idx, layout_name, hint in sample_configs:
            if layout_idx >= len(layouts):
                continue
            layout = layouts[layout_idx]
            slide = self.prs.slides.add_slide(layout)
            self._fill_sample_placeholders(slide, layout_name, hint)

    def _fill_sample_placeholders(self, slide: Any, layout_name: str, hint: str) -> None:
        """在样例 slide 的 placeholders 中填入提示文本。"""
        for shape in slide.placeholders:
            ph_type = shape.placeholder_format.type
            if ph_type == 1:  # TITLE
                shape.text_frame.text = f"[{layout_name}] 单击此处添加标题"
                self._set_placeholder_font(shape, "微软雅黑", 28, self.TEXT_GRAY)
            elif ph_type in (2, 7):  # BODY or OBJECT
                shape.text_frame.text = f"{hint}\n\n• 要点一\n• 要点二\n• 要点三"
                self._set_placeholder_font(shape, "微软雅黑", 18, self.TEXT_GRAY)
            elif ph_type == 3:  # CENTER_TITLE
                shape.text_frame.text = "单击此处添加标题"
                self._set_placeholder_font(shape, "微软雅黑", 36, self.TEXT_GRAY)
            elif ph_type == 4:  # SUBTITLE
                shape.text_frame.text = "单击此处添加副标题"
                self._set_placeholder_font(shape, "微软雅黑", 20, self.TEXT_GRAY)

    def _set_placeholder_font(
        self,
        shape: Any,
        font_name: str,
        font_size: int,
        font_color: RGBColor,
    ) -> None:
        """设置 placeholder 中所有文本的字体。"""
        if not shape.has_text_frame:
            return
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.name = font_name
                run.font.size = Pt(font_size)
                run.font.color.rgb = font_color


# ========== CLI 入口 ==========

def main(argv: Optional[List[str]] = None) -> int:
    parser = argparse.ArgumentParser(description="课件模板构建器")
    parser.add_argument("--output", required=True, type=Path, help="输出模板路径(.pptx)")
    parser.add_argument("--verbose", action="store_true", help="详细日志")
    args = parser.parse_args(argv)

    if args.verbose:
        logging.getLogger().setLevel(logging.DEBUG)

    try:
        builder = TemplateBuilder()
        builder.build(args.output)
        print(f"✅ 模板创建成功: {args.output}")
        return 0
    except Exception as exc:
        logger.error("构建失败: %s", exc)
        print(f"❌ 构建失败: {exc}")
        return 1


if __name__ == "__main__":
    sys.exit(main())
